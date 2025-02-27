import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_standard_presentation():
    # Load the custom template with the orange background already set up.
    prs = Presentation("custom_template.pptx")
    
    # Function: Add a slide with a title, bullet points and optionally an image.
    def add_bullet_slide(title_text, bullet_points, img_path=None, img_pos=(Inches(5), Inches(1.5)), img_width=Inches(3)):
        slide_layout = prs.slide_layouts[5]  # Use a blank custom layout from the template.
        slide = prs.slides.add_slide(slide_layout)
        
        # Title textbox
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4.5), Inches(1))
        title_tf = title_box.text_frame
        title_tf.text = title_text
        for p in title_tf.paragraphs:
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT
        
        # Bullet points textbox
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4))
        bullet_tf = bullet_box.text_frame
        
        if bullet_points:
            bullet_tf.text = bullet_points[0]
            bullet_tf.paragraphs[0].font.size = Pt(20)
            bullet_tf.paragraphs[0].font.name = "Calibri"
            bullet_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            for bp in bullet_points[1:]:
                p = bullet_tf.add_paragraph()
                p.text = bp
                p.level = 1
                p.font.size = Pt(20)
                p.font.name = "Calibri"
                p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Optionally add an image (logo or illustration)
        if img_path:
            try:
                slide.shapes.add_picture(img_path, img_pos[0], img_pos[1], width=img_width)
            except Exception as e:
                print(f"Error adding image '{img_path}': {e}")
        
        return slide

    # Function: Add a two-column slide (for “Before” vs. “After” content).
    def add_two_column_slide(title_text, left_title, left_bullets, right_title, right_bullets):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Slide title (centered)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.text = title_text
        for p in title_tf.paragraphs:
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Left column (Before)
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
        left_tf = left_box.text_frame
        left_tf.text = left_title
        for p in left_tf.paragraphs:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(255, 255, 255)
        for bullet in left_bullets:
            p = left_tf.add_paragraph()
            p.text = bullet
            p.level = 1
            p.font.size = Pt(20)
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Right column (After)
        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
        right_tf = right_box.text_frame
        right_tf.text = right_title
        for p in right_tf.paragraphs:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(255, 255, 255)
        for bullet in right_bullets:
            p = right_tf.add_paragraph()
            p.text = bullet
            p.level = 1
            p.font.size = Pt(20)
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        return slide

    # Function: Add a long bullet list slide (for GTM Strategy details).
    def add_long_bullet_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title textbox
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.text = title_text
        for p in title_tf.paragraphs:
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Bullet points textbox
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5))
        bullet_tf = bullet_box.text_frame
        
        if bullet_points:
            bullet_tf.text = bullet_points[0]
            bullet_tf.paragraphs[0].font.size = Pt(20)
            bullet_tf.paragraphs[0].font.name = "Calibri"
            bullet_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            for bp in bullet_points[1:]:
                p = bullet_tf.add_paragraph()
                p.text = bp
                p.level = 1
                p.font.size = Pt(20)
                p.font.name = "Calibri"
                p.font.color.rgb = RGBColor(255, 255, 255)
        
        return slide

    # ----------------------------
    # Slide 1: Why Valuezen? (Problem Statement & Opportunity)
    slide1_title = "Why Valuezen? (Problem Statement & Opportunity)"
    slide1_bullets = [
        "Decision-making in logistics is slow due to lack of real-time, quantifiable value insights.",
        "Mid-market and SMBs struggle with cost justification, integration complexity, and slow adoption.",
        "Existing platforms offer data but lack personalized ROI-driven intelligence."
    ]
    # Add the first slide with the logo (assumes "images/logo.png" is in place)
    add_bullet_slide(slide1_title, slide1_bullets,
                     img_path="images/logo.png",
                     img_pos=(Inches(7), Inches(0.3)), img_width=Inches(2))

    # ----------------------------
    # Slide 2: What is Valuezen? (Solution Overview)
    slide2_title = "What is Valuezen? (Solution Overview)"
    slide2_bullets = [
        "AI-powered value delivery platform for logistics & transportation.",
        "Plug-and-play API-first approach for rapid deployment.",
        "Live value calculators to showcase impact in cost, time, and efficiency."
    ]
    add_bullet_slide(slide2_title, slide2_bullets,
                     img_path="images/solution_image.png",
                     img_pos=(Inches(5), Inches(1.5)), img_width=Inches(3))

    # ----------------------------
    # Slide 3: Key Benefits (Before vs. After Valuezen)
    slide3_title = "Key Benefits (Before vs. After Valuezen)"
    left_title = "Before:"
    left_bullets = [
        "Fragmented data, unclear ROI.",
        "Slow customer onboarding & adoption.",
        "Lack of AI-driven decision intelligence."
    ]
    right_title = "After:"
    right_bullets = [
        "API-first selling → faster deployment, proven cost savings.",
        "AI/ML-driven value insights → predictive decision-making.",
        "Free-tier trials → SMB adoption & viral expansion."
    ]
    add_two_column_slide(slide3_title, left_title, left_bullets, right_title, right_bullets)

    # ----------------------------
    # Slide 4: GTM Strategy & Expansion Plan
    slide4_title = "GTM Strategy & Expansion Plan"
    slide4_bullets = [
        "Target Mid-Market with Rapid Deployment & API Integration",
        "   • Plug-and-play solution with API-first selling approach.",
        "   • Faster response times → reduced downtime & automation-led savings.",
        "   • Build case studies to show mid-market efficiency gains.",
        "",
        "Drive SMB Adoption with Free Tier & Simplified Onboarding",
        "   • Simplified UX → highlight ease of use & time savings in marketing.",
        "   • Free-tier tracking service for small carriers & brokers.",
        "   • Leverage referrals & integrate with popular TMS platforms.",
        "",
        "Expand into Latin America & APAC Through Mid-Market Penetration",
        "   • Localized content, multilingual support, and regional partnerships.",
        "   • Flexible pricing to match emerging market needs.",
        "",
        "Strengthen Competitive Differentiation with AI & Predictive Insights",
        "   • AI/ML-powered predictive ETA & automation as differentiators.",
        "   • Target C-level executives with data-driven supply chain intelligence.",
        "   • Position Valuezen as a decision intelligence leader."
    ]
    add_long_bullet_slide(slide4_title, slide4_bullets)

    # Save the presentation to a BytesIO buffer.
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

if __name__ == "__main__":
    ppt_buffer = create_standard_presentation()
    with open("Valuezen_Standard_Presentation.pptx", "wb") as f:
        f.write(ppt_buffer.getbuffer())
    print("Presentation generated as 'Valuezen_Standard_Presentation.pptx'.")
